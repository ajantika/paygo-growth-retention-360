"""Generate branded PPTX decks from in-app data.

Two exports:
- build_exec_summary_deck(kpis, mrr_trend)        Overview page
- build_account_deck(account_id, ctx)             Account 360 page

Uses python-pptx native shapes/charts so it works on Streamlit Cloud
without external binaries.
"""
from __future__ import annotations

import io

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ORANGE = RGBColor(0xF3, 0x80, 0x20)
NAVY = RGBColor(0x00, 0x36, 0x82)
GREY = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(12); sp.font.color.rgb = WHITE


def _add_kpi_card(slide, x: float, y: float, w: float, h: float, label: str, value: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GREY
    card.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(h - 0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(10); p1.font.color.rgb = GREY

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(22); p2.font.bold = True; p2.font.color.rgb = NAVY


def _add_footer(slide, text: str = "Synthetic demo data — seed=42 — Ajantika Paul") -> None:
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9); p.font.color.rgb = GREY; p.font.italic = True


def _add_line_chart(slide, x: float, y: float, w: float, h: float, df: pd.DataFrame,
                    x_col: str, y_col: str, title: str) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [pd.Timestamp(d).strftime("%b %Y") for d in df[x_col]]
    chart_data.add_series(y_col, df[y_col].tolist())
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    )
    chart = chart_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.has_legend = False


# ---------- Exec summary deck ----------

def build_exec_summary_deck(kpis: dict, mrr_trend: pd.DataFrame) -> bytes:
    """1 slide: KPI cockpit + total-MRR trend."""
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    _add_title_bar(
        s, "PayGo Growth & Retention 360 — Exec Summary",
        subtitle=f"As of {kpis['latest_month'].strftime('%B %Y')}  •  Demo data",
    )

    # KPI cards row
    cards = [
        ("Current MRR", _money(kpis["current_mrr"])),
        ("Active accounts", f"{kpis['active_accounts']:,}"),
        ("NRR", f"{kpis['nrr']:.1f}%"),
        ("GRR", f"{kpis['grr']:.1f}%"),
        ("Graduation rate", f"{kpis['graduation_rate']:.1f}%"),
        ("Median time-to-upgrade", f"{int(kpis['median_time_to_upgrade_days'])} days"),
    ]
    w = 2.0; gap = 0.15; total = len(cards) * w + (len(cards) - 1) * gap
    start_x = (13.33 - total) / 2
    for i, (label, value) in enumerate(cards):
        _add_kpi_card(s, start_x + i * (w + gap), 1.3, w, 1.0, label, value)

    _add_line_chart(s, 1.0, 2.8, 11.3, 3.7, mrr_trend, "month", "mrr", "Total MRR over time")

    _add_footer(s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------- Account deck ----------

def build_account_deck(account_id: str, ctx: dict, account_row: pd.Series) -> bytes:
    """3-slide per-account drilldown: title, snapshot, MRR trajectory."""
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    mrr = ctx["mrr"]
    products = ctx["products"]
    churn = ctx["churn"]
    grad = ctx["graduation"]

    # Slide 1: title
    s1 = prs.slides.add_slide(blank)
    _add_title_bar(s1, f"Account 360 — {account_id}", subtitle="Per-account growth & retention report")
    tb = s1.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        f"Region: {account_row['region']}",
        f"Acquisition channel: {account_row['channel']}",
        f"Plan type: {account_row['plan_type']}",
        f"Entry product: {account_row['entry_product']}",
        f"Signup: {pd.Timestamp(account_row['signup_month']).strftime('%b %Y')}",
        f"Current segment: {account_row['current_segment']}",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18); p.font.color.rgb = NAVY
    _add_footer(s1)

    # Slide 2: snapshot KPIs + key events
    s2 = prs.slides.add_slide(blank)
    _add_title_bar(s2, f"Snapshot — {account_id}")

    tenure_months = 0
    if not mrr.empty:
        first_m = pd.Timestamp(mrr["month"].min())
        last_m = pd.Timestamp(mrr["month"].max())
        tenure_months = (last_m.year - first_m.year) * 12 + (last_m.month - first_m.month) + 1
    latest_mrr = float(mrr["mrr"].iloc[-1]) if not mrr.empty else 0.0
    peak_mrr = float(mrr["mrr"].max()) if not mrr.empty else 0.0
    n_products = int(products["product"].nunique()) if not products.empty else 0

    snapshot = [
        ("Tenure", f"{tenure_months} months"),
        ("Latest MRR", _money(latest_mrr)),
        ("Peak MRR", _money(peak_mrr)),
        ("Products adopted", f"{n_products}"),
    ]
    w = 2.7; gap = 0.2
    for i, (label, value) in enumerate(snapshot):
        _add_kpi_card(s2, 0.6 + i * (w + gap), 1.4, w, 1.1, label, value)

    # Key events list
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key events"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY

    events: list[tuple[pd.Timestamp, str]] = []
    if not products.empty:
        for _, pr in products.iterrows():
            events.append((pd.Timestamp(pr["start_month"]), f"Adopted {pr['product']}"))
    if not grad.empty:
        g = grad.iloc[0]
        events.append((pd.Timestamp(g["graduation_month"]),
                       f"Graduated to Enterprise (time-to-upgrade: {int(g['time_to_upgrade_days'])} days)"))
    if not churn.empty:
        c = churn.iloc[0]
        events.append((pd.Timestamp(c["churn_month"]),
                       f"Churned — reason: {c['churn_reason']}"))
    events.sort(key=lambda e: e[0])
    for dt, label in events:
        ep = tf.add_paragraph()
        ep.text = f"  •  {dt.strftime('%b %Y')} — {label}"
        ep.font.size = Pt(12); ep.font.color.rgb = NAVY

    _add_footer(s2)

    # Slide 3: MRR trajectory
    s3 = prs.slides.add_slide(blank)
    _add_title_bar(s3, f"MRR trajectory — {account_id}")
    if not mrr.empty:
        _add_line_chart(s3, 0.6, 1.2, 12.1, 5.5, mrr.assign(month=pd.to_datetime(mrr["month"])),
                        "month", "mrr", "Monthly recurring revenue")
    _add_footer(s3)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _money(x: float) -> str:
    if abs(x) >= 1_000_000:
        return f"${x/1_000_000:.2f}M"
    if abs(x) >= 1_000:
        return f"${x/1_000:.1f}K"
    return f"${x:,.0f}"
